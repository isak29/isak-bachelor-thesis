from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── Knightec colour palette (light theme) ────────────────────────
BG        = RGBColor(0xFF, 0xFF, 0xFF)   # white background
CARD      = RGBColor(0xF4, 0xF5, 0xF7)   # light card
CARD_DK   = RGBColor(0xE8, 0xE9, 0xEB)   # slightly darker card
ORANGE    = RGBColor(0xF2, 0x65, 0x22)   # Knightec orange
ORANGE_DK = RGBColor(0xBF, 0x4A, 0x10)   # dark orange
TEXT      = RGBColor(0x1A, 0x1A, 0x1A)   # near-black main text
GRAY      = RGBColor(0x55, 0x55, 0x55)   # medium gray
GRAY_LT   = RGBColor(0x99, 0x99, 0x99)   # light gray
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
BLUE      = RGBColor(0x2E, 0x86, 0xC1)

# ── Helpers ───────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    _rect(slide, 0, 0, 13.33, 0.1, ORANGE, None)
    _rect(slide, 0, 7.4, 13.33, 0.1, CARD, None)
    return slide

def _rect(slide, l, t, w, h, fill, line):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=14, bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def heading(slide, text):
    tb(slide, text, 0.5, 0.15, 12.33, 0.72, size=26, bold=True, color=TEXT)
    _rect(slide, 0.5, 0.95, 12.33, 0.04, ORANGE, None)

def img_box(slide, l, t, w, h, label="[ BILD ]"):
    _rect(slide, l, t, w, h, CARD, GRAY_LT)
    tb(slide, label, l, t + h/2 - 0.25, w, 0.5, size=13, italic=True,
       color=GRAY, align=PP_ALIGN.CENTER)

def bullet(slide, items, l, t, w, size=13):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, col) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = col
    return box

# ── Presentation setup ────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – Titel
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)

_rect(s, 0, 2.8, 13.33, 0.06, ORANGE, None)

tb(s, "Effektivisering av informationssökning med LLM:er",
   0.8, 0.9, 11.73, 1.7, size=34, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

tb(s, "En jämförelse av grafbaserad kunskapsgraf mot API-baserad informationshämtning",
   1.0, 3.05, 11.33, 0.7, size=16, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

tb(s, "Isak Lampell  ·  Knightec Group AB  ·  Mittuniversitetet  ·  2026",
   1.0, 4.1, 11.33, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – Problemet (övergripande)
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Problemet")

tb(s, "Föreställ dig att du sitter på ett företag med hundratals kollegor,\nmånga projekt och tusentals dokumentationssidor med information om riktlinjer och rutiner.",
   0.5, 1.15, 8.5, 0.9, size=15, color=TEXT)

tb(s, "Allt är spritt över många olika interna system och du behöver snabbt hitta svaret på en specifik fråga.",
   0.5, 2.2, 8.5, 0.65, size=15, color=TEXT)

tb(s, "Det du gör är att du börjar söka, frågar kollegor och söker lite ytterligare\ntills du hittar svaret. Det kan ta väldigt lång tid som egentligen borde\nläggas på det riktiga arbetet.",
   0.5, 3.05, 8.5, 1.1, size=15, color=GRAY)

_rect(s, 0.5, 4.45, 8.5, 1.3, CARD, ORANGE)
tb(s, "Det här är vardagen hos många organisationer idag —\noch det är precis det problemet mitt arbete syftar till att lösa.",
   0.65, 4.6, 8.2, 1.0, size=15, bold=True, color=TEXT)

img_box(s, 9.2, 1.15, 3.8, 5.6, "[ BILD ]")

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – Bakgrund: Knightec & Systemen
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Bakgrund – Knightec")

tb(s, "Exjobbet är gjort i samarbete med Knightec — ett tech-konsultbolag med verksamhet på flera orter i Sverige.",
   0.5, 1.1, 12.3, 0.55, size=14, color=TEXT)

tb(s, "Informationen är spridd över tre system:", 0.5, 1.8, 8.0, 0.45, size=14, bold=True, color=ORANGE)

systems = [
    ("GitHub", "Projekt & programkod\nRepositories, commits,\nissues, aktiva användare", BLUE),
    ("Slack", "Kommunikation\nKanaler, meddelanden,\nanvändare", GREEN),
    ("Confluence", "Intern dokumentation\nProcesser, rutiner,\nriktlinjer, teams", ORANGE),
]
for i, (name, desc, col) in enumerate(systems):
    x = 0.5 + i * 4.2
    _rect(s, x, 2.4, 3.9, 3.9, CARD, col)
    tb(s, name, x, 2.5, 3.9, 0.65, size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    _rect(s, x + 0.3, 3.25, 3.3, 0.04, col, None)
    tb(s, desc, x + 0.15, 3.4, 3.6, 2.6, size=13, color=GRAY, align=PP_ALIGN.CENTER)

tb(s, "Problemet: ingen gemensam struktur — data kan inte enkelt kombineras",
   0.5, 6.45, 12.3, 0.5, size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – Teori
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Teori")

boxes = [
    ("LLM", "Large Language Model — AI-modell tränad på stora mängder text.\n\nI arbetet: tolkar användarens fråga och avgör vilka verktyg som ska anropas.", BLUE),
    ("MCP — Model Context Protocol", "Mellanlager mellan LLM och datakällor. Funktionalitet exponeras som verktyg.\n\nI arbetet: MCP-server med verktyg för Cypher-frågor eller API-anrop.", ORANGE),
    ("Kunskapsgraf", "Nätverksstruktur av entiteter och relationer. Lagras i Neo4j.\n\nI arbetet: data från GitHub, Slack och Confluence modellerat med noder och kanter.", GREEN),
    ("Cypher", "Frågespråk för Neo4j. LLM:en genererar Cypher-frågor dynamiskt.\n\nEx: MATCH (p:Person)-[:WORKS_ON]->(r:Repository) RETURN r", GRAY),
]
positions = [(0.4, 1.15), (6.9, 1.15), (0.4, 4.0), (6.9, 4.0)]
for (lbl, desc, col), (x, y) in zip(boxes, positions):
    _rect(s, x, y, 5.9, 2.65, CARD, col)
    tb(s, lbl, x + 0.15, y + 0.1, 5.6, 0.5, size=17, bold=True, color=col)
    _rect(s, x + 0.15, y + 0.65, 5.6, 0.03, col, None)
    tb(s, desc, x + 0.15, y + 0.8, 5.6, 1.7, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – Kunskapsgraf (bild)
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Kunskapsgraf – visualisering")

tb(s, "Grafen innehåller noder och relationer från alla tre system. Samma person är kopplad till repos, Slack-inlägg och Confluence-dokument.",
   0.5, 1.1, 12.3, 0.55, size=14, italic=True, color=ORANGE)

img_box(s, 0.5, 1.8, 8.5, 5.35, "[ BILD: visualisation.png — Neo4j grafvisualisering ]")

bullet(s, [
    ("Noder:", TEXT),
    ("Person, Repository, Issue, Commit,", GRAY),
    ("SlackChannel, Post, Document ...", GRAY),
    ("", TEXT),
    ("Relationer:", TEXT),
    ("WORKS_ON  ·  COMMITTED_IN", GRAY),
    ("AUTHORED  ·  MEMBER_OF", GRAY),
    ("HAS_POST  ·  WROTE ...", GRAY),
], l=9.2, t=2.0, w=3.9, size=12)

_rect(s, 9.2, 5.5, 3.9, 1.65, CARD, ORANGE)
tb(s, "Tydliga kopplingar gör att LLM:en kan följa sambanden direkt — utan att gissa.",
   9.3, 5.65, 3.7, 1.35, size=13, bold=True, color=TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – Metod
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Metod — Design Science Research")

tb(s, "Designa och bygg konkreta artefakter — utvärdera dem för att dra slutsatser.",
   0.5, 1.1, 12.3, 0.45, size=14, italic=True, color=ORANGE)

steps_dsr = ["Designa\ntvå lösningar", "→", "Implementera\nGraf + API", "→", "Utvärdera\n15 testfrågor", "→", "Dra\nslutstatser"]
xpos = [0.4, 2.65, 3.05, 5.3, 5.7, 7.95, 8.35]
widths = [2.2, 0.35, 2.2, 0.35, 2.2, 0.35, 2.2]
for i, (txt, x, w) in enumerate(zip(steps_dsr, xpos, widths)):
    if txt == "→":
        tb(s, "→", x, 2.0, w, 0.7, size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    else:
        _rect(s, x, 1.75, w, 1.1, CARD, ORANGE)
        tb(s, txt, x + 0.05, 1.85, w - 0.1, 0.9, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

levels = [
    ("Simpla  Q1–Q5", "Data från en enskild källa.\nInget komplext resonemang.", BLUE),
    ("Medelsvåra  Q6–Q10", "Mer resonemang krävs.\nFlera anrop från en källa.", ORANGE),
    ("Komplexa  Q11–Q15", "Kopplingar mellan\nolika datakällor krävs.", RED),
]
for i, (title, desc, col) in enumerate(levels):
    x = 0.4 + i * 4.25
    _rect(s, x, 3.2, 4.05, 3.95, CARD, col)
    tb(s, title, x + 0.1, 3.3, 3.85, 0.65, size=15, bold=True, color=col)
    _rect(s, x + 0.3, 4.05, 3.45, 0.04, col, None)
    tb(s, desc, x + 0.15, 4.2, 3.75, 1.3, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, "Mäts: korrekthet · anrop · tid",
       x + 0.15, 5.7, 3.75, 0.4, size=11, italic=True, color=GRAY_LT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – Systemarkitektur
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Systemarkitektur")

arch_boxes = [
    (0.4,  3.1, 2.0, 1.0, "Användare",               CARD, ORANGE),
    (2.9,  3.1, 2.4, 1.0, "Cherry Studio\n(UI + chatt)", CARD, ORANGE),
    (5.9,  1.7, 2.4, 1.0, "LLM",                     CARD, BLUE),
    (5.9,  4.2, 2.4, 1.0, "MCP Server",              CARD, ORANGE),
    (9.2,  1.7, 3.9, 1.0, "Lösning 1: Graf\n(Neo4j)",       CARD, GREEN),
    (9.2,  4.2, 3.9, 1.0, "Lösning 2: API\nGitHub/Slack/Confluence", CARD, RED),
]
for (x, y, w, h, lbl, fill, border) in arch_boxes:
    _rect(s, x, y, w, h, fill, border)
    tb(s, lbl, x + 0.05, y + 0.08, w - 0.1, h - 0.16, size=13, bold=True,
       color=TEXT, align=PP_ALIGN.CENTER)

tb(s, "→", 2.3, 3.45, 0.65, 0.5, size=22, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "↑↓", 6.85, 2.1, 0.5, 0.5, size=18, color=BLUE, align=PP_ALIGN.CENTER)
tb(s, "↑↓", 6.85, 4.6, 0.5, 0.5, size=18, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "Cypher →", 8.35, 2.05, 1.1, 0.35, size=11, italic=True, color=GREEN)
tb(s, "REST →",   8.35, 4.6, 1.1, 0.35, size=11, italic=True, color=RED)

_rect(s, 0.4, 5.9, 12.5, 1.25, CARD, GRAY_LT)
tb(s, "1. Fråga ställs  →  2. LLM väljer verktyg  →  3. MCP kör verktyget  →  4. Data returneras  →  5. LLM formulerar svar",
   0.55, 6.05, 12.2, 0.9, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – Grafbaserade lösningen
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Lösning 1 — Grafbaserad")

_rect(s, 0.5, 1.15, 6.0, 5.9, CARD, ORANGE)
tb(s, "Hur det fungerar", 0.65, 1.25, 5.7, 0.5, size=17, bold=True, color=ORANGE)

steps = [
    "1.  LLM hämtar grafens schema\n    (alla noder, relationer, attribut)",
    "2.  LLM genererar en Cypher-fråga\n    baserat på schemat och frågan",
    "3.  Frågan körs mot Neo4j",
    "4.  Resultat returneras — klart",
]
for i, step in enumerate(steps):
    y = 1.9 + i * 1.1
    _rect(s, 0.65, y, 0.45, 0.45, ORANGE, None)
    tb(s, str(i+1), 0.65, y, 0.45, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, 1.2, y, 5.1, 0.85, size=13, color=TEXT)

_rect(s, 0.65, 6.35, 5.7, 0.55, WHITE, GREEN)
tb(s, "Totalt: 2 verktyg  (get_schema + execute_query)",
   0.7, 6.4, 5.6, 0.45, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

img_box(s, 6.8, 1.15, 6.2, 5.9, "[ BILD: Exempelflöde /\nNeo4j visualisering ]")

# ═══════════════════════════════════════════════════════════
# SLIDE 9 – API-baserade lösningen
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Lösning 2 — API-baserad")

tb(s, "Ingen central lagring — data hämtas i realtid direkt från systemens API:er",
   0.5, 1.1, 12.3, 0.45, size=14, italic=True, color=ORANGE)

api_systems = [
    ("GitHub", "8 verktyg", BLUE, ["get_repositories", "get_commits", "get_issues",
                                    "get_active_users", "get_org_members", "..."]),
    ("Slack", "5 verktyg", GREEN, ["get_channels", "get_messages",
                                    "get_users", "search_messages", "..."]),
    ("Confluence", "7 verktyg", ORANGE, ["get_spaces", "get_pages",
                                          "get_page_content", "search_content", "..."]),
]
for i, (name, count, col, tools) in enumerate(api_systems):
    x = 0.4 + i * 4.3
    _rect(s, x, 1.7, 4.0, 0.65, col, None)
    tb(s, f"{name}  —  {count}", x, 1.75, 4.0, 0.55, size=15, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    for j, tool in enumerate(tools):
        _rect(s, x, 2.45 + j * 0.58, 4.0, 0.52, CARD, CARD_DK)
        tb(s, f"  {tool}", x + 0.1, 2.5 + j * 0.58, 3.8, 0.42, size=12, color=GRAY)

_rect(s, 1.9, 6.55, 9.5, 0.65, CARD, ORANGE)
tb(s, "Totalt: 20 verktyg   vs   2 verktyg i den grafbaserade lösningen",
   1.9, 6.6, 9.5, 0.55, size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 – Demo
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Demo")

tb(s, "Förinspelad video — grafbaserade lösningen i aktion",
   0.5, 1.1, 12.3, 0.45, size=16, italic=True, color=ORANGE)

img_box(s, 1.5, 1.7, 10.33, 5.3, "[ VIDEO ]")

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – Resultat: Korrekthet
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Resultat — Korrekthet")

cd = ChartData()
cd.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd.add_series('Grafbaserad (%)', (100, 80, 80))
cd.add_series('API-baserad (%)',  (100, 40, 20))

ch_shape = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.9), cd)
ch = ch_shape.chart
ch.has_legend = True
ch.legend.position = 2
ch.legend.include_in_layout = False
plot = ch.plots[0]
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = ORANGE
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = BLUE
ch.value_axis.maximum_scale = 100
ch.value_axis.minimum_scale = 0

obs = [
    (GREEN,  "Simpla:",       "Båda 100% ✓"),
    (ORANGE, "Medelsvåra:",   "Graf 80%  —  API 40%"),
    (RED,    "Komplexa:",     "Graf 80%  —  API 20%"),
]
for i, (col, label, val) in enumerate(obs):
    y = 1.3 + i * 1.3
    _rect(s, 8.2, y, 4.9, 1.1, CARD, col)
    tb(s, label, 8.35, y + 0.1, 4.6, 0.4, size=14, bold=True, color=col)
    tb(s, val,   8.35, y + 0.55, 4.6, 0.45, size=16, bold=True, color=TEXT)

_rect(s, 8.2, 5.3, 4.9, 1.8, CARD, RED)
tb(s, "API:s problem:", 8.35, 5.4, 4.6, 0.4, size=14, bold=True, color=RED)
tb(s, "Kan ej länka identiteter\nmellan system.\n\"isak29\" i GitHub ≠\n\"Isak Lampell\" i Slack?",
   8.35, 5.85, 4.6, 1.15, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 – Resultat: Verktygsanrop & Svarstid
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Resultat — Verktygsanrop & Svarstid")

# Left: verktygsanrop
tb(s, "Genomsnittliga verktygsanrop", 0.4, 1.1, 6.3, 0.35, size=12, bold=True, color=GRAY)
cd2 = ChartData()
cd2.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd2.add_series('Grafbaserad', (2.2, 4.0, 2.4))
cd2.add_series('API-baserad', (1.6, 4.4, 6.0))
ch2 = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.5), Inches(6.1), Inches(5.6), cd2).chart
ch2.has_legend = True
ch2.legend.position = 2
ch2.legend.include_in_layout = False
plot2 = ch2.plots[0]
plot2.series[0].format.fill.solid(); plot2.series[0].format.fill.fore_color.rgb = ORANGE
plot2.series[1].format.fill.solid(); plot2.series[1].format.fill.fore_color.rgb = BLUE
ch2.value_axis.maximum_scale = 8

# Right: svarstid
tb(s, "Genomsnittlig svarstid (ms)", 6.8, 1.1, 6.3, 0.35, size=12, bold=True, color=GRAY)
cd3 = ChartData()
cd3.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd3.add_series('Grafbaserad', (1967, 5234, 2965))
cd3.add_series('API-baserad', (1521, 5251, 18388))
ch3 = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.6), cd3).chart
ch3.has_legend = True
ch3.legend.position = 2
ch3.legend.include_in_layout = False
plot3 = ch3.plots[0]
plot3.series[0].format.fill.solid(); plot3.series[0].format.fill.fore_color.rgb = ORANGE
plot3.series[1].format.fill.solid(); plot3.series[1].format.fill.fore_color.rgb = BLUE
ch3.value_axis.maximum_scale = 20000

# divider
_rect(s, 6.66, 1.1, 0.04, 5.9, GRAY_LT, None)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 – Diskussion & Slutsats
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Diskussion & Slutsats")

tb(s, "Vad visar resultaten?", 0.5, 1.1, 6.0, 0.45, size=16, bold=True, color=ORANGE)
bullet(s, [
    ("+ Graf ger LLM:en tydliga relationer att följa", GREEN),
    ("", TEXT),
    ("+ En Cypher-fråga kan navigera hela grafen", GREEN),
    ("", TEXT),
    ("− Schema-overhead vid varje anrop", RED),
    ("", TEXT),
    ("− API kan ej länka identiteter mellan system", RED),
    ("", TEXT),
    ("+ API fungerar bra vid enkla frågor", GREEN),
], l=0.7, t=1.7, w=5.8, size=13)

_rect(s, 6.9, 1.1, 6.1, 4.4, CARD, ORANGE)
tb(s, "Slutsats", 7.05, 1.2, 5.8, 0.5, size=18, bold=True, color=ORANGE)
tb(s,
   "En grafbaserad representation ger LLM:en tydliga fördelar när information behöver kombineras från olika källor.\n\nDen ger LLM:en tydliga relationer som är enkla att följa.\n\nAPI-lösningen fungerar bra för frågor mot enskilda system.",
   7.05, 1.85, 5.8, 3.4, size=14, color=TEXT)

_rect(s, 6.9, 5.7, 6.1, 1.45, CARD, GREEN)
tb(s, "Framtida arbete", 7.05, 5.8, 5.8, 0.4, size=14, bold=True, color=GREEN)
tb(s, "Kombinera lösningarna: API bygger grafen, grafen används för sökning",
   7.05, 6.25, 5.8, 0.75, size=13, italic=True, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 14 – Tack & Frågor
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)

_rect(s, 0, 3.3, 13.33, 0.06, ORANGE, None)

tb(s, "Tack!", 1.5, 1.0, 10.33, 2.0, size=80, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
tb(s, "Frågor?", 1.5, 3.6, 10.33, 1.2, size=42, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

tb(s, "Isak Lampell  ·  lampell.i@gmail.com\nKnightec Group AB  ·  Mittuniversitetet  ·  2026",
   1.5, 5.4, 10.33, 0.85, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "thesis_presentation_1.pptx")
prs.save(out)
print(f"Saved: {out}")
