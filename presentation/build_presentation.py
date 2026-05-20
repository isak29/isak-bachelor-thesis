from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu
import os

# ── Colour palette ────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
CARD_BG    = RGBColor(0x16, 0x27, 0x3A)
TEAL       = RGBColor(0x00, 0xB4, 0xD8)
DEEP_BLUE  = RGBColor(0x00, 0x77, 0xB6)
LIGHT_TEAL = RGBColor(0x90, 0xE0, 0xEF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xA0, 0xB4, 0xC0)
GREEN      = RGBColor(0x00, 0xD4, 0xAA)
ORANGE     = RGBColor(0xFF, 0xAA, 0x00)
RED        = RGBColor(0xFF, 0x5C, 0x5C)

# ── Helpers ───────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    # top bar
    _rect(slide, 0, 0, 13.33, 0.07, TEAL, None)
    # bottom bar
    _rect(slide, 0, 7.43, 13.33, 0.07, DEEP_BLUE, None)
    return slide

def _rect(slide, l, t, w, h, fill, line):
    from pptx.util import Inches
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def tb_lines(slide, lines, l, t, w, h, size=13, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, line_spacing=None):
    """Multi-paragraph textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bld
        run.font.color.rgb = clr
    return box

def hline(slide, y, color=TEAL, x0=0.5, x1=12.83):
    from pptx.util import Inches, Pt
    conn = slide.shapes.add_connector(1, Inches(x0), Inches(y), Inches(x1), Inches(y))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.2)

def heading(slide, text, y=0.18):
    tb(slide, text, 0.5, y, 12.3, 0.75, size=28, bold=True, color=WHITE)
    hline(slide, 1.06)

def card(slide, l, t, w, h, border_color=TEAL):
    _rect(slide, l, t, w, h, CARD_BG, border_color)

# ── Build ─────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)

tb(s, "Effektivisering av informationssökning med LLM:er",
   0.7, 1.6, 11.9, 2.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "En jämförelse av grafbaserad kunskapsgraf mot API-baserad informationshämtning",
   1.0, 3.7, 11.3, 0.8, size=18, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

_rect(s, 3.8, 4.85, 5.73, 0.05, TEAL, None)

tb(s, "Isak Lampell  ·  Knightec Group AB  ·  Mittuniversitetet  ·  2026",
   1.0, 5.15, 11.3, 0.55, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – Bakgrund & Problem
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Bakgrund & Problem")

# Left column – text
tb(s, "Datasilos i organisationer", 0.5, 1.2, 6.2, 0.45, size=17, bold=True, color=TEAL)
tb(s, "• Data är utspridd i separata system\n• Svårt att få en samlad bild av verksamheten\n• Manuell sökning är tidskrävande", 0.5, 1.7, 6.2, 1.1, size=14, color=WHITE)

tb(s, "Kostnaden är hög", 0.5, 2.95, 6.2, 0.45, size=17, bold=True, color=TEAL)
tb(s, "• En kunskapsarbetare spenderar i genomsnitt 2,5h/dag\n  att söka information  (IDC-studie)\n• Det motsvarar ~30% av en arbetsdag\n• Globalt förlorar företag miljarder på grund av detta", 0.5, 3.45, 6.2, 1.3, size=14, color=WHITE)

tb(s, "AI:s potential blockeras", 0.5, 4.9, 6.2, 0.45, size=17, bold=True, color=TEAL)
tb(s, "• AI är beroende av sammanhängande, kontextuell data\n• Datasilos gör att många AI-initiativ stannar vid prototyper", 0.5, 5.4, 6.2, 0.9, size=14, color=WHITE)

# Right – stat cards
card(s, 7.3, 1.2, 5.6, 2.4, TEAL)
tb(s, "2,5 h", 7.3, 1.3, 5.6, 1.35, size=64, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, "per dag söker anställda information", 7.3, 2.75, 5.6, 0.6, size=14, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

card(s, 7.3, 3.85, 5.6, 2.65, TEAL)
tb(s, "Möjlig lösning", 7.4, 3.95, 5.4, 0.45, size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, "Kunskapsgrafer representerar data som\nnoder och relationer – ett sammanhängande\nnätverk av organisationens information.\n\nLLM:er kan navigera grafen för att\nsnabbt hitta och kombinera relevant data.",
   7.4, 4.5, 5.4, 1.8, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – Syfte & Verifierbara mål
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Syfte & Verifierbara mål")

tb(s, "Undersöka hur organisatorisk information kan representeras för att effektivisera "
      "informationssökning med hjälp av LLM:er – genom att jämföra en grafbaserad struktur "
      "mot en API-baserad lösning.",
   0.5, 1.15, 12.3, 0.75, size=14, italic=True, color=LIGHT_TEAL)

goals = [
    ("1", "Effektivitet", "Hur lång tid tar det för en LLM\natt hitta och hämta information\ni grafbaserat system vs API?"),
    ("2", "Korrekthet",   "Hur ofta genereras relevanta\noch korrekta svar?"),
    ("3", "Resursanvändning", "Hur många verktygsanrop\nbehövs för att hitta\ninformationen?"),
]
for i, (num, title, desc) in enumerate(goals):
    x = 0.4 + i * 4.3
    card(s, x, 2.1, 4.1, 4.7, TEAL)
    tb(s, num, x, 2.15, 4.1, 1.3, size=60, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _rect(s, x + 0.3, 3.55, 3.5, 0.05, TEAL, None)
    tb(s, title, x + 0.1, 3.7, 3.9, 0.55, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc,  x + 0.1, 4.35, 3.9, 2.1, size=13, color=GRAY, align=PP_ALIGN.CENTER)

tb(s, "Genomförs i samarbete med Knightec Group AB",
   0.5, 6.9, 12.3, 0.4, size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – Systemarkitektur
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Systemarkitektur")

# Flow boxes
def flow_box(slide, x, y, w, h, label, fill, border):
    _rect(slide, x, y, w, h, fill, border)
    tb(slide, label, x+0.05, y+0.08, w-0.1, h-0.16, size=13, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

flow_box(s, 0.3,  3.0, 2.0, 1.1, "Användare",             CARD_BG, TEAL)
flow_box(s, 3.0,  3.0, 2.4, 1.1, "Cherry Studio\n(UI + Agent)", CARD_BG, TEAL)
flow_box(s, 5.9,  1.7, 2.4, 1.1, "LLM\n(GPT-4o-mini)",    RGBColor(0x0C,0x28,0x40), DEEP_BLUE)
flow_box(s, 5.9,  4.3, 2.4, 1.1, "MCP Server\n(TypeScript)", RGBColor(0x00,0x2A,0x38), TEAL)
flow_box(s, 9.2,  1.7, 3.8, 1.1, "Lösning 1\nGrafdatabas (Neo4j)",         RGBColor(0x06,0x22,0x18), GREEN)
flow_box(s, 9.2,  4.3, 3.8, 1.1, "Lösning 2\nGitHub + Slack + Confluence", RGBColor(0x22,0x08,0x08), RED)

# Arrow labels
tb(s, "→ fråga",             2.3, 3.35, 0.75, 0.4, size=10, color=GRAY, italic=True)
tb(s, "sysp. + verktyg →",   4.75, 2.1, 1.5, 0.35, size=9, color=GRAY, italic=True)
tb(s, "← svar",             4.75, 2.55, 1.5, 0.3, size=9, color=GRAY, italic=True)
tb(s, "HTTP POST →",         4.75, 4.55, 1.5, 0.35, size=9, color=GRAY, italic=True)
tb(s, "← JSON",             4.75, 4.95, 1.5, 0.3, size=9, color=GRAY, italic=True)
tb(s, "Cypher →",            8.35, 2.1, 0.9, 0.3, size=9, color=GREEN, italic=True)
tb(s, "REST →",              8.35, 4.55, 0.9, 0.3, size=9, color=RED, italic=True)

# Step legend
card(s, 0.3, 5.55, 12.7, 1.65, DEEP_BLUE)
tb(s, "Flöde per fråga:", 0.5, 5.6, 3.0, 0.4, size=13, bold=True, color=TEAL)
tb(s, "1.  Användaren ställer fråga i Cherry Studio\n"
      "2.  Cherry skickar systemprompten + verktygslista + chatthistorik till LLM:en\n"
      "3.  LLM:en beslutar om verktygsanrop → Cherry skickar HTTP POST till MCP-servern\n"
      "4.  MCP returnerar JSON → LLM kan kedja flera anrop (agent-loop) → svar till användaren",
   0.5, 6.05, 12.4, 1.1, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – Lösning 1: Grafbaserad
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Lösning 1: Grafbaserad (Knowledge Graph)")

# Left
tb(s, "Datamodell", 0.5, 1.15, 5.8, 0.45, size=17, bold=True, color=TEAL)
tb(s, "Data från GitHub, Slack och Confluence lagras\nsom noder och relationer i Neo4j.\n\n"
      "Noder: Person · Team · Project · Repository\n"
      "          Issue · Commit · SlackChannel · Post\n"
      "          ConfluenceSpace · Document\n\n"
      "Relationer:  MEMBER_OF · AUTHORED · HAS_POST\n"
      "                  COMMITTED_IN · HAS_ISSUE  …",
   0.5, 1.65, 5.8, 3.1, size=13, color=WHITE)

tb(s, "Verktyg (2 st)", 0.5, 4.85, 5.8, 0.45, size=17, bold=True, color=TEAL)
_rect(s, 0.5, 5.37, 5.8, 0.55, CARD_BG, TEAL)
tb(s, "get_schema  — hämtar hela grafens schema",        0.6, 5.43, 5.5, 0.43, size=13, color=WHITE)
_rect(s, 0.5, 6.0, 5.8, 0.55, CARD_BG, TEAL)
tb(s, "execute_query  — kör Cypher-fråga mot Neo4j", 0.6, 6.06, 5.5, 0.43, size=13, color=WHITE)

# Right – mini graph illustration
tb(s, "Exempelflöde: vem commitade i payment-service?",
   6.8, 1.15, 6.2, 0.45, size=13, bold=True, color=TEAL)

nodes = [
    (7.1,  2.1, 1.9, 0.75, ":Person\nIsak Lampell",   TEAL),
    (9.6,  1.6, 2.2, 0.75, ":Repository\npayment-service", GREEN),
    (9.6,  3.3, 2.2, 0.75, ":Document\nAWS-info",     ORANGE),
    (7.1,  3.8, 1.9, 0.75, ":Commit\nabc123",          LIGHT_TEAL),
]
for (nx, ny, nw, nh, label, col) in nodes:
    _rect(s, nx, ny, nw, nh, CARD_BG, col)
    tb(s, label, nx+0.06, ny+0.06, nw-0.12, nh-0.12, size=10, color=col, align=PP_ALIGN.CENTER)

rel_labels = [
    (8.95, 2.15, "COMMITTED_IN →"),
    (8.5,  3.5,  "← AUTHORED"),
]
for (rx, ry, rl) in rel_labels:
    tb(s, rl, rx, ry, 2.2, 0.35, size=9, italic=True, color=GRAY)

tb(s, "LLM:en hämtar all denna data\nmed en enda Cypher-fråga\n(efter schema-anrop).",
   6.8, 5.0, 6.2, 0.9, size=13, italic=True, color=GRAY)

_rect(s, 6.8, 6.05, 6.2, 0.65, CARD_BG, GREEN)
tb(s, "Fördel: ett traverseringsanrop ersätter\nmånga separata API-anrop",
   6.9, 6.1, 6.0, 0.55, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – Lösning 2: API-baserad
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Lösning 2: API-baserad")

tb(s, "Ingen central lagring — data hämtas dynamiskt via varje systems egna API:er",
   0.5, 1.15, 12.3, 0.45, size=14, italic=True, color=LIGHT_TEAL)

services = [
    ("GitHub",     "8 verktyg", GREEN,  [
        "get_repositories", "get_issues", "get_commits",
        "get_active_users", "get_repo_info", "get_org_members",
        "get_commit_detail", "search_issues"]),
    ("Slack",      "5 verktyg", ORANGE, [
        "get_channels", "get_messages",
        "get_users", "search_messages", "get_channel_history"]),
    ("Confluence", "7 verktyg", TEAL,   [
        "get_spaces", "get_pages", "get_page_content",
        "search_content", "get_authors", "get_page_by_id",
        "get_recent_pages"]),
]
for i, (name, count, col, tools) in enumerate(services):
    x = 0.35 + i * 4.35
    _rect(s, x, 1.75, 4.15, 0.62, col, None)
    tb(s, f"{name}  ({count})", x, 1.8, 4.15, 0.52, size=16, bold=True,
       color=DARK_BG, align=PP_ALIGN.CENTER)
    for j, tool in enumerate(tools):
        ty = 2.45 + j * 0.52
        _rect(s, x, ty, 4.15, 0.46, CARD_BG, RGBColor(0x28,0x40,0x50))
        tb(s, f"  {tool}", x, ty+0.06, 4.15, 0.34, size=11, color=GRAY)

_rect(s, 3.3, 6.8, 6.73, 0.5, CARD_BG, TEAL)
tb(s, "Totalt: 20 verktyg  vs  2 verktyg i den grafbaserade lösningen",
   3.3, 6.82, 6.73, 0.45, size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – Evalueringsupplägg
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Evalueringsupplägg")

tb(s, "15 testfrågor · Varje fråga körs i ny session · Mäter: korrekthet, svarstid, verktygsanrop",
   0.5, 1.15, 12.3, 0.45, size=14, italic=True, color=LIGHT_TEAL)

levels = [
    ("Simpla",     "Q1–Q5",   TEAL,   "Hämtar data från en enda\nkälla. Inget resonemang\neller kopplingar krävs.",
     '"Get all GitHub repos"\n"Show latest Slack message"'),
    ("Medelsvåra", "Q6–Q10",  ORANGE, "Flera anrop inom samma\nkälla, eller ytterligare\nresonemang krävs.",
     '"What GitHub issues are open\nin payment-service?"'),
    ("Komplexa",   "Q11–Q15", RED,    "Tydliga kopplingar mellan\nolika datakällor krävs.\nKomplexa resonemang.",
     '"Name of latest Confluence doc\nfrom person who did last commit\nin payment-service?"'),
]
for i, (title, qrange, col, desc, ex) in enumerate(levels):
    x = 0.35 + i * 4.35
    _rect(s, x, 1.78, 4.15, 5.25, CARD_BG, col)
    tb(s, title,  x, 1.85, 4.15, 0.7, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s, qrange, x, 2.6,  4.15, 0.45, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    _rect(s, x+0.4, 3.15, 3.35, 0.05, col, None)
    tb(s, desc, x+0.1, 3.3, 3.95, 1.35, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(s, x+0.1, 4.75, 3.95, 0.05, RGBColor(0x30,0x45,0x55), None)
    tb(s, "Exempel:", x+0.15, 4.85, 3.85, 0.35, size=11, bold=True, color=GRAY)
    tb(s, ex, x+0.15, 5.25, 3.85, 1.5, size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – Resultat: Korrekthet
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Resultat: Korrekthet")

cd = ChartData()
cd.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd.add_series('Grafbaserad (%)', (100, 80, 80))
cd.add_series('API-baserad (%)',  (100, 40, 20))

ch_shape = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.25), Inches(7.3), Inches(5.7), cd)
ch = ch_shape.chart
ch.has_legend = True
ch.legend.position = 2
ch.legend.include_in_layout = False
plot = ch.plots[0]
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = TEAL
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = RED
ch.value_axis.maximum_scale = 100
ch.value_axis.minimum_scale = 0

# Right panel
tb(s, "Viktigaste observationer", 8.0, 1.25, 5.0, 0.45, size=17, bold=True, color=TEAL)

obs = [
    (GREEN,  True,  "Simpla frågor:",   "Båda 100% korrekta"),
    (ORANGE, False, "Medelsvåra:",       "Graf 80%  —  API 40%"),
    (RED,    False, "Komplexa:",         "Graf 80%  —  API 20%"),
    (GRAY,   False, "API:s problem:",    "Kan inte länka identiteter\nmellan olika system\n(ex. isak29 i GitHub =\nIsak Lampell i Slack?)"),
    (TEAL,   False, "Graf behåller:",    "Konstant 80% korrekthet\növer alla komplexitetsnivåer"),
]
y = 1.85
for col, bld, label, detail in obs:
    _rect(s, 7.9, y, 0.15, 0.35, col, None)
    tb(s, label,  8.15, y, 2.1,  0.4, size=13, bold=True, color=col)
    tb(s, detail, 10.2, y, 2.85, 0.7, size=13, color=WHITE)
    y += 0.9

# ═══════════════════════════════════════════════════════════
# SLIDE 9 – Resultat: Svarstid & Verktygsanrop
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Resultat: Svarstid & Verktygsanrop")

# Time chart
cd2 = ChartData()
cd2.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd2.add_series('Grafbaserad (ms)', (1967, 5234, 2965))
cd2.add_series('API-baserad (ms)', (1521, 5251, 18388))

ch2 = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.3), Inches(6.5), Inches(3.5), cd2).chart
ch2.has_legend = True
ch2.legend.position = 2
plot2 = ch2.plots[0]
plot2.series[0].format.fill.solid(); plot2.series[0].format.fill.fore_color.rgb = TEAL
plot2.series[1].format.fill.solid(); plot2.series[1].format.fill.fore_color.rgb = RED
ch2.value_axis.maximum_scale = 20000

tb(s, "Genomsnittlig svarstid (ms)", 0.4, 1.1, 7.0, 0.35, size=12, italic=True, color=GRAY)

# Tool-calls chart
cd3 = ChartData()
cd3.categories = ['Simpla', 'Medelsvåra', 'Komplexa']
cd3.add_series('Grafbaserad', (2.2, 4.0, 2.4))
cd3.add_series('API-baserad', (1.6, 4.4, 6.0))

ch3 = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(5.1), Inches(6.5), Inches(2.05), cd3).chart
ch3.has_legend = False
plot3 = ch3.plots[0]
plot3.series[0].format.fill.solid(); plot3.series[0].format.fill.fore_color.rgb = TEAL
plot3.series[1].format.fill.solid(); plot3.series[1].format.fill.fore_color.rgb = RED
ch3.value_axis.maximum_scale = 8

tb(s, "Genomsnittliga verktygsanrop", 0.4, 4.9, 7.0, 0.35, size=12, italic=True, color=GRAY)

# Right: key numbers
tb(s, "Nyckelsiffror", 7.3, 1.3, 5.7, 0.45, size=17, bold=True, color=TEAL)

stats = [
    ("Komplexa — svarstid",  "Graf: 2 965 ms",  "API: 18 388 ms", TEAL, RED),
    ("Komplexa — anrop",     "Graf: snitt 2,4", "API: snitt 6,0",  TEAL, RED),
    ("Q12 (extremfall)",     "Graf: 2 993 ms",  "API: 72 314 ms(!)", TEAL, RED),
]
for i, (label, v1, v2, c1, c2) in enumerate(stats):
    y0 = 1.9 + i * 1.5
    card(s, 7.3, y0, 5.7, 1.35, DEEP_BLUE)
    tb(s, label, 7.45, y0+0.05, 5.4, 0.4, size=12, bold=True, color=GRAY)
    tb(s, v1, 7.45, y0+0.55, 2.7, 0.55, size=16, bold=True, color=c1)
    tb(s, v2, 10.1, y0+0.55, 2.7, 0.55, size=16, bold=True, color=c2)

_rect(s, 7.3, 6.55, 5.7, 0.7, CARD_BG, TEAL)
tb(s, "Graf är effektivast vid komplexa frågor.\nAPI är snabbare och enklare vid enkla.",
   7.4, 6.6, 5.5, 0.6, size=13, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 – Diskussion
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Diskussion")

# Graf column
tb(s, "Grafbaserade lösningen", 0.5, 1.15, 5.8, 0.45, size=17, bold=True, color=TEAL)
graph_items = [
    (GREEN, "+ Tydliga kopplingar möjliggör traversering\n   mellan datakällor"),
    (GREEN, "+ Lågt och stabilt antal verktygsanrop"),
    (GREEN, "+ Flexibel sökning via genererade Cypher-frågor"),
    (RED,   "− Schema-overhead vid varje fråga"),
    (RED,   "− Kräver manuell uppbyggnad och underhåll av grafen"),
]
for i, (col, text) in enumerate(graph_items):
    y0 = 1.7 + i * 0.72
    _rect(s, 0.5, y0+0.05, 0.14, 0.34, col, None)
    tb(s, text, 0.75, y0, 5.4, 0.55, size=13, color=WHITE)

# API column
tb(s, "API-baserade lösningen", 7.0, 1.15, 5.8, 0.45, size=17, bold=True, color=RED)
api_items = [
    (GREEN, "+ Snabb och resurssnål vid enkla frågor"),
    (GREEN, "+ Alltid aktuell data – ingen lagring krävs"),
    (RED,   "− Kan inte länka identiteter mellan system"),
    (RED,   "− Fler verktygsanrop vid komplexa frågor"),
    (RED,   "− Korrekthet sjunker kraftigt vid komplexa frågor"),
]
for i, (col, text) in enumerate(api_items):
    y0 = 1.7 + i * 0.72
    _rect(s, 7.0, y0+0.05, 0.14, 0.34, col, None)
    tb(s, text, 7.25, y0, 5.6, 0.55, size=13, color=WHITE)

# Methodology box
card(s, 0.4, 5.45, 12.5, 1.75, DEEP_BLUE)
tb(s, "Metodreflektioner", 0.6, 5.5, 4.0, 0.4, size=14, bold=True, color=TEAL)
tb(s, "• Frågorna designades av mig — risk att de gynnar grafbaserade lösningen\n"
      "• Varje fråga kördes en gång — slumpen kan påverka enskilda resultat\n"
      "• Korrekthetsbedömning gjordes subjektivt — i gränsfall är LLM:ens resonemang viktigare än svaret",
   0.6, 5.95, 12.1, 1.0, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – Slutsats & Framtida arbete
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)
heading(s, "Slutsats")

conclusions = [
    (TEAL,   "Grafbaserad struktur ger LLM:en bättre förutsättningar att navigera och kombinera\n"
              "information från flera datakällor."),
    (GREEN,  "Komplexa frågor: Graf 80% korrekt på ~3 000 ms  —  API 20% korrekt på ~18 000 ms"),
    (ORANGE, "API-lösningen presterar väl vid enkla frågor men tappar kraftigt vid kopplingar\nmellan system."),
    (LIGHT_TEAL, "Resultaten stödjer tidigare forskning: kunskapsgrafer ger LLM:er bättre\nstrukturerad kontext och förbättrar svarskvaliteten."),
]
y = 1.25
for col, text in conclusions:
    _rect(s, 0.4, y+0.1, 0.14, 0.4, col, None)
    tb(s, text, 0.65, y, 12.3, 0.7, size=14, color=WHITE)
    y += 0.95

_rect(s, 0.4, 5.35, 12.5, 0.05, TEAL, None)
tb(s, "Framtida arbete", 0.5, 5.55, 5.0, 0.45, size=17, bold=True, color=TEAL)
future = [
    "• Automatisera processen för att bygga och uppdatera grafen (API → Graf pipeline)",
    "• Kombinera de två lösningarna: välj approach dynamiskt utifrån frågans komplexitet",
    "• Testa lösningen med verklig organisatorisk data i produktion",
]
for i, f in enumerate(future):
    tb(s, f, 0.5, 6.05 + i * 0.42, 12.3, 0.4, size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 – Tack & Frågor
# ═══════════════════════════════════════════════════════════
s = new_slide(prs)

tb(s, "Tack!", 1.5, 1.4, 10.33, 1.8, size=80, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

_rect(s, 3.8, 3.55, 5.73, 0.06, TEAL, None)

tb(s, "Frågor?", 1.5, 3.8, 10.33, 1.0, size=40, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "Isak Lampell  ·  lampell.i@gmail.com\nKnightec Group AB  ·  Mittuniversitetet  ·  2026",
   1.5, 5.3, 10.33, 0.85, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "thesis_presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
